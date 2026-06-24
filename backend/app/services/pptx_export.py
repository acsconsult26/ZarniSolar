from __future__ import annotations
"""Export pipeline: merge a project's data into template.pptx and produce
the final downloadable deck. Tokens are replaced inside existing runs only
(never rebuilds paragraphs), so original formatting is preserved. Burmese
runs are forced onto a Myanmar-capable font so they never render as boxes."""
import io
import re
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

from ..schema import merged_field_values
from .flowchart import render_priority_flowchart

NAVY = RGBColor(0x0D, 0x2C, 0x54)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x2A, 0x35)

# category -> (slide index 1-based, spec table shape name, title shape name)
SPEC_SLIDES = {
    "inverter": (14, "Table 1", "TextBox 3"),
    "battery": (15, "Table 4", "TextBox 7"),
    "panel": (16, "Table 1", "Rectangle 2"),
}

TEMPLATE_PATH = Path(__file__).resolve().parent.parent / "template.pptx"

TOKEN_RE = re.compile(r"\{\{(\w+)\}\}")
BURMESE_RE = re.compile(r"[က-႟]")
BURMESE_FONT = "Pyidaungsu"

# slide index (1-based) -> shape name -> field key in `uploads` dict
IMAGE_SLOTS = {
    3: {"Picture 13": "satellite_photo"},
    4: {
        "Picture 1": "survey_photo_1",
        "Picture 2": "survey_photo_2",
        "Picture 3": "survey_photo_3",
        "Picture 5": "survey_photo_4",
        "Picture 4": "survey_photo_5",
        "Picture 8": "survey_photo_6",
    },
    5: {
        "Picture 1": "data_logger_photo_1",
        "Picture 2": "data_logger_photo_2",
        "Picture 4": "data_logger_photo_3",
    },
    6: {"Picture 20": "load_profile_chart"},
    9: {"Image": "system_drawing"},
}


def find_shape(shapes, name):
    for shape in shapes:
        if shape.shape_type == 6:
            found = find_shape(shape.shapes, name)
            if found is not None:
                return found
        if shape.name == name:
            return shape
    return None


def _text_frames_of(shape):
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                yield cell.text_frame
    elif shape.has_text_frame:
        yield shape.text_frame


def replace_tokens(prs, values: dict):
    for slide in prs.slides:
        for shape in slide.shapes:
            _replace_in_shape(shape, values)


def _replace_in_shape(shape, values: dict):
    if shape.shape_type == 6:
        for sub in shape.shapes:
            _replace_in_shape(sub, values)
        return
    for tf in _text_frames_of(shape):
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                if "{{" not in run.text:
                    continue
                new_text = TOKEN_RE.sub(lambda m: str(values.get(m.group(1), "")), run.text)
                if new_text != run.text:
                    run.text = new_text
                if BURMESE_RE.search(run.text):
                    run.font.name = BURMESE_FONT


def _replace_picture(slide, shape_name: str, image_bytes: bytes):
    shape = find_shape(slide.shapes, shape_name)
    if shape is None:
        return False
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    shape._element.getparent().remove(shape._element)
    slide.shapes.add_picture(io.BytesIO(image_bytes), left, top, width, height)
    return True


def _replace_picture_contain(slide, shape_name: str, image_bytes: bytes):
    """Swap a placeholder picture but preserve the new image's aspect ratio,
    fitting it inside the placeholder box (no stretching) and centering it."""
    shape = find_shape(slide.shapes, shape_name)
    if shape is None:
        return False
    box_left, box_top, box_w, box_h = shape.left, shape.top, shape.width, shape.height
    shape._element.getparent().remove(shape._element)
    try:
        with Image.open(io.BytesIO(image_bytes)) as im:
            iw, ih = im.size
    except Exception:
        iw, ih = box_w, box_h
    scale = min(box_w / iw, box_h / ih)
    new_w = int(iw * scale)
    new_h = int(ih * scale)
    left = box_left + (box_w - new_w) // 2
    top = box_top + (box_h - new_h) // 2
    slide.shapes.add_picture(io.BytesIO(image_bytes), left, top, new_w, new_h)
    return True


def _placeholder_image(text: str = "No image added") -> bytes:
    """A neutral light-grey placeholder PNG used when no photo was uploaded,
    so leftover reference-deck images never appear in a client's proposal."""
    img = Image.new("RGB", (800, 600), (235, 238, 242))
    draw = ImageDraw.Draw(img)
    draw.rectangle([2, 2, 797, 597], outline=(180, 188, 200), width=3)
    # default bitmap font is tiny; scale up by drawing then noting it's fine for a placeholder
    bbox = draw.textbbox((0, 0), text)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    draw.text(((800 - tw) / 2, (600 - th) / 2), text, fill=(120, 130, 145))
    out = io.BytesIO()
    img.save(out, format="PNG")
    return out.getvalue()


def _set_shape_text(shape, text: str):
    """Replace a shape's text with `text`, keeping the first run's font but
    forcing a Myanmar font for Burmese. Paragraphs split on newlines."""
    if shape is None or not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.word_wrap = True
    # capture a font size from the first existing run if any
    size = None
    for p in tf.paragraphs:
        for r in p.runs:
            if r.font.size:
                size = r.font.size
                break
        if size:
            break
    tf.clear()
    lines = text.split("\n")
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = line
        if BURMESE_RE.search(line):
            run.font.name = BURMESE_FONT
        run.font.size = size or Pt(12)


def _set_first_run_text(shape, text: str):
    """Set a title shape's text while keeping the first run's formatting."""
    if shape is None or not shape.has_text_frame:
        return
    paras = shape.text_frame.paragraphs
    if not paras or not paras[0].runs:
        return
    paras[0].runs[0].text = text
    for r in paras[0].runs[1:]:
        r.text = ""
    if BURMESE_RE.search(text):
        paras[0].runs[0].font.name = BURMESE_FONT


def _build_spec_table(slide, table_shape_name: str, product: dict):
    """Replace the brand-specific template table with a clean 3-column
    (Specification / Value / Unit) table generated from the chosen product's
    structured specs. Keeps the original table's position and size."""
    specs = product.get("specs") or []
    if not specs:
        return  # nothing to generate; leave template table untouched
    shape = find_shape(slide.shapes, table_shape_name)
    if shape is None:
        return
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    shape._element.getparent().remove(shape._element)

    rows = len(specs) + 1
    gf = slide.shapes.add_table(rows, 3, left, top, width, height)
    table = gf.table
    table.columns[0].width = int(width * 0.55)
    table.columns[1].width = int(width * 0.30)
    table.columns[2].width = int(width * 0.15)

    def style_cell(cell, text, *, header=False):
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if header else WHITE
        tf = cell.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = ""
        run = tf.paragraphs[0].add_run()
        run.text = str(text)
        run.font.size = Pt(11) if header else Pt(10)
        run.font.bold = header
        run.font.color.rgb = WHITE if header else DARK
        if BURMESE_RE.search(str(text)):
            run.font.name = BURMESE_FONT

    style_cell(table.cell(0, 0), product.get("spec_title") or "Specification", header=True)
    style_cell(table.cell(0, 1), "Value", header=True)
    style_cell(table.cell(0, 2), "Unit", header=True)
    for i, spec in enumerate(specs, start=1):
        style_cell(table.cell(i, 0), spec.get("label", ""))
        style_cell(table.cell(i, 1), spec.get("value", ""))
        style_cell(table.cell(i, 2), spec.get("unit", ""))


def _apply_selected_products(prs, selected_products: dict):
    for category, (slide_idx, table_name, title_name) in SPEC_SLIDES.items():
        product = selected_products.get(category)
        if not product:
            continue
        slide = prs.slides[slide_idx - 1]
        title = product.get("spec_title") or f"{product.get('brand', '')} {product.get('model_name', '')}".strip()
        if title:
            _set_first_run_text(find_shape(slide.shapes, title_name), title)
        _build_spec_table(slide, table_name, product)

    # Slide 22 warranty - product warranty lines from the chosen products
    lines = []
    for category in ("panel", "inverter", "battery"):
        product = selected_products.get(category)
        if product and (product.get("warranty_line") or "").strip():
            lines.append(product["warranty_line"].strip())
    if lines:
        slide22 = prs.slides[21]
        _set_shape_text(find_shape(slide22.shapes, "Rectangle 47"), "\n".join(lines))


def _load_image_bytes(uploads: dict, field: str, storage):
    path = uploads.get(field)
    if not path or not storage.exists(path):
        return None
    return storage.read_bytes(path)


def export_project(project, storage, reference_images: list[bytes] | None = None,
                   selected_products: dict | None = None) -> bytes:
    """project: has .data (dict), .uploads (dict), .slide19_image_path, .flowchart_image_path
    selected_products: {category: product_dict} for slides 14-16 spec tables + warranty"""
    values = merged_field_values(project.data or {})
    prs = Presentation(str(TEMPLATE_PATH))
    replace_tokens(prs, values)

    if selected_products:
        _apply_selected_products(prs, selected_products)

    uploads = project.uploads or {}
    for slide_idx, slots in IMAGE_SLOTS.items():
        slide = prs.slides[slide_idx - 1]
        for shape_name, field in slots.items():
            img = _load_image_bytes(uploads, field, storage)
            # Always swap the slot: an uploaded photo, or a neutral "No image
            # added" placeholder so the original reference-deck photos never
            # leak into a client's proposal.
            _replace_picture(slide, shape_name, img or _placeholder_image())

    # Slide 18 - Power Management: per-client free text overrides the template
    pm_text = (values.get("power_management_text") or "").strip()
    if pm_text:
        slide18 = prs.slides[17]
        _set_shape_text(find_shape(slide18.shapes, "Rectangle 1"), pm_text)
        for extra in ("Rectangle 2", "Rectangle 3"):
            sh = find_shape(slide18.shapes, extra)
            if sh is not None:
                sh._element.getparent().remove(sh._element)

    # Slide 21 - Power Source Priority: edited/auto-drafted text overrides template
    pp_text = (values.get("power_priority_text") or "").strip()
    if pp_text:
        slide21 = prs.slides[20]
        _set_shape_text(find_shape(slide21.shapes, "Rectangle 3"), pp_text)

    # Slide 19 - AI-generated infographic (cached image on the project, or manual fallback)
    if storage.exists(project.slide19_image_path):
        img19 = storage.read_bytes(project.slide19_image_path)
        slide19 = prs.slides[18]
        left = top = Emu(0)
        width = prs.slide_width
        height = prs.slide_height
        slide19.shapes.add_picture(io.BytesIO(img19), left, top, width, height)

    # Slide 20 - data-driven flowchart (drawn, not AI), inserted over the placeholder
    flowchart_path = project.flowchart_image_path
    if not (flowchart_path and storage.exists(flowchart_path)):
        flowchart_bytes = render_priority_flowchart(values)
    else:
        flowchart_bytes = storage.read_bytes(flowchart_path)
    slide20 = prs.slides[19]
    _replace_picture_contain(slide20, "Picture 3", flowchart_bytes)

    # Slide 23 - admin reference gallery (replace as many photo slots as we have images for)
    if reference_images:
        slide23 = prs.slides[22]
        photo_shapes = [
            s for s in slide23.shapes
            if s.shape_type == 13 and s.name.startswith("Picture")
        ]
        for shape, img in zip(photo_shapes, reference_images):
            _replace_picture(slide23, shape.name, img)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
