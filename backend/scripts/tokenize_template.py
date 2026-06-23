"""
One-time build step: read the original company template and write a tokenized
copy (template.pptx) used by the export pipeline at runtime.

Run: python backend/scripts/tokenize_template.py

This performs *substring* replacement inside existing runs (it never rebuilds
paragraphs), so all original run-level formatting (font, size, color, bold,
Burmese vs Latin font) is preserved untouched.
"""
import sys
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[2]
SRC = ROOT / "MRTV Prosal 17.June.2026.pptx"
DST = ROOT / "backend" / "app" / "template.pptx"

# (slide_index 1-based, shape_name, old_substring, token)
# Order matters: longer/more specific substrings first to avoid partial clobber.
REPLACEMENTS = [
    # Slide 1 - cover
    (1, "TextBox 10", "MRTV", "{{site_name}}"),
    (1, "TextBox 15", "17.June .2026", "{{proposal_date}}"),
    (1, "TextBox 16", "09-2031977", "{{contact_phone}}"),

    # Slide 3 - project background
    (3, "TextBox 6", "MRTV", "{{site_name}}"),
    (3, "TextBox 6", "ဆွမ်းဦးပုံညရှင် )", "{{site_name_mm}} )"),
    (3, "TextBox 6", "21.9022473, 95.992391", "{{gps_lat}}, {{gps_lng}}"),
    (3, "TextBox 7", "110 kVA", "{{generator_capacity_kva}} kVA"),
    (3, "TextBox 14", "88.59 kW", "{{max_load_kw}} kW"),
    (3, "TextBox 14", "-  48 kW", "-  {{avg_load_kw}} kW"),
    (3, "TextBox 14", "  -  6h ", "  -  {{min_grid_supply}}"),
    (3, "TextBox 14", "~", ""),
    (3, "TextBox 14", " 8h ", ""),

    # Slide 5 - data logger
    (5, "Rectangle 3", "(22-May-2026 , 10:35AM)", "({{data_logger_start}})"),
    (5, "Rectangle 3", "(24-May-2026 , 10:35AM)", "({{data_logger_end}})"),

    # Slide 6 - load profile
    (6, "TextBox 1", "88.59kWp ", "{{max_load_kw}}kWp "),
    (6, "TextBox 1", "= 5kWp", "= {{min_load_kw}}kWp"),
    (6, "TextBox 1", "= 48kWp", "= {{avg_load_kw}}kWp"),
    (6, "TextBox 1", "230V", "{{voltage_v}}V"),
    (6, "TextBox 1", "Power Factor   = ", "Power Factor   = {{power_factor}}"),
    (6, "Rectangle 2", "900 Units", "{{daily_usage_units}} Units"),

    # Slide 17 - system technical configuration table (cells are split across runs)
    (17, "Table 2", "Sigenstack", "{{battery_module}}"),
    (17, "Table 2", "Sigen", "{{inverter_model}}"),
    (17, "Table 2", " 60kW M1  HYB (x2)", " (x{{inverter_qty}})"),
    (17, "Table 2", "120 kW Total Output", "{{total_inverter_kw}} kW Total Output"),
    (17, "Table 2", " 12.0kW", ""),
    (17, "Table 2", " (x36)", " (x{{battery_qty}})"),
    (17, "Table 2", "432 kWh", "{{total_battery_kwh}} kWh"),
    (17, "Table 2", "Longi", "{{panel_brand}}"),
    (17, "Table 2", " 650kW Mono", " {{panel_watt}}W Mono"),
    (17, "Table 2", " (x288)", " (x{{panel_qty}})"),
    (17, "Table 2", "187.2 ", "{{total_solar_kwp}} "),
    (17, "Table 2", "8", "{{backup_hours}}"),
    (17, "Table 2", "Design Margin : 20% Reserve ", "Design Margin : {{design_margin_pct}}% Reserve "),

    # Slide 18 - power management / savings
    (18, "Rectangle 3", " 432kWh ", " {{total_battery_kwh}}kWh "),
    (18, "Rectangle 3", "၇၀% မှ ၈၀%", "{{bill_savings_low_pct}}% မှ {{bill_savings_high_pct}}%"),

    # Slide 21 - power source priority text
    (21, "Rectangle 3", "Load 48kW ", "Load {{solar_load_kw}}kW "),
    (21, "Rectangle 3", "432 kWh ", "{{total_battery_kwh}} kWh "),
    (21, "Rectangle 3", "20% ", "{{generator_dod_trigger_pct}}% "),

    # Slide 22 - warranty
    (22, "Rectangle 43", "1 Year) Warranty ", "{{install_warranty_years}} Year) Warranty "),
]

# Slide 8 has a couple of values split across awkward runs (Burmese digit
# grouping). For these only, match the *exact* original run text rather than
# a substring, and clear the trailing sibling run instead of leaving stale text.
# (slide_index, shape_name, exact_old_run_text, token_for_that_run, clear_next_n_runs)
EXACT_RUN_REPLACEMENTS = [
    (8, "Rectangle 18", "၂", "{{panel_qty}}", 1),   # '၂' + '၀၀ ' -> "200 ..." (panel count)
    (8, "Rectangle 18", "1", "{{total_solar_kwp}}", 1),  # '1' + '68 ' -> "168 kWp" (AUTO total)
]


def apply_exact_run_replacements(prs):
    for slide_idx, shape_name, old, token, clear_next in EXACT_RUN_REPLACEMENTS:
        slide = prs.slides[slide_idx - 1]
        shape = find_shape(slide.shapes, shape_name)
        if shape is None or not shape.has_text_frame:
            print(f"  ! shape not found: slide {slide_idx} / {shape_name}")
            continue
        runs = [r for p in shape.text_frame.paragraphs for r in p.runs]
        for i, run in enumerate(runs):
            if run.text == old:
                run.text = token
                for j in range(1, clear_next + 1):
                    if i + j < len(runs):
                        runs[i + j].text = ""
                break
        else:
            print(f"  ! exact run not found: slide {slide_idx} / {shape_name} / {old!r}")


# Image placeholder shapes that get swapped at export time, mapped by
# (slide_index, shape_name) -> field name. Detected by inspecting the original
# deck; the shape named 'Image' on content slides is the swappable photo.
IMAGE_SHAPES = {
    (3, "Picture 13"): "satellite_photo",
    (4, "Picture 1"): "survey_photo_1",
    (4, "Picture 2"): "survey_photo_2",
    (4, "Picture 3"): "survey_photo_3",
    (4, "Picture 4"): "survey_photo_4",
    (4, "Picture 5"): "survey_photo_5",
    (4, "Picture 8"): "survey_photo_6",
    (5, "Picture 1"): "data_logger_photo_1",
    (5, "Picture 2"): "data_logger_photo_2",
    (5, "Picture 4"): "data_logger_photo_3",
    (6, "Picture 20"): "load_profile_chart",
    (9, "Image"): "system_drawing",  # placeholder slot, may stay fixed if not uploaded
    (19, "AutoShape 4"): "ai_infographic",  # slide 19 body swapped wholesale by service
    (20, "Picture 3"): "priority_flowchart",
}


def find_shape(shapes, name):
    for shape in shapes:
        if shape.shape_type == 6:  # GROUP
            found = find_shape(shape.shapes, name)
            if found is not None:
                return found
        if shape.name == name:
            return shape
    return None


def _text_frames_of(shape):
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                yield cell.text_frame
    elif shape.has_text_frame:
        yield shape.text_frame


def apply_text_replacements(prs):
    for slide_idx, shape_name, old, token in REPLACEMENTS:
        slide = prs.slides[slide_idx - 1]
        shape = find_shape(slide.shapes, shape_name)
        if shape is None or not (shape.has_text_frame or shape.has_table):
            print(f"  ! shape not found or no text frame: slide {slide_idx} / {shape_name}")
            continue
        replaced = False
        for tf in _text_frames_of(shape):
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    if old and old in run.text:
                        run.text = run.text.replace(old, token)
                        replaced = True
        if not replaced:
            print(f"  ! substring not found: slide {slide_idx} / {shape_name} / {old!r}")


def main():
    if not SRC.exists():
        print(f"Source template not found at {SRC}", file=sys.stderr)
        sys.exit(1)
    prs = Presentation(str(SRC))
    apply_text_replacements(prs)
    apply_exact_run_replacements(prs)
    DST.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(DST))
    print(f"Wrote tokenized template to {DST}")


if __name__ == "__main__":
    main()
