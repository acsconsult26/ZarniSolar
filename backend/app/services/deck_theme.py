"""Dark-themed 16:9 deck engine for the redesigned Zarni proposal.

All slides use a clean dark background with subtle technical pattern shapes and
white text. A professional font hierarchy (title / subtitle / body / footer) is
applied via helpers. Burmese runs are forced onto a Myanmar font.
"""
from __future__ import annotations

import re
import threading

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ---- palettes (4 selectable deck templates) ----
# WHITE/MUTED/RED stay constant across templates -- text legibility and the
# error/accent-red meaning shouldn't change with the color scheme.
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0xB8, 0xC6, 0xDA)
RED = RGBColor(0xE2, 0x23, 0x1A)

THEMES = {
    "navy": {
        "label": "Navy & Gold", "description": "The original Zarni look — deep navy with a gold accent.",
        "BG": RGBColor(0x0B, 0x14, 0x22), "PANEL": RGBColor(0x13, 0x22, 0x38),
        "LINE": RGBColor(0x1B, 0x30, 0x4C), "ACCENT": RGBColor(0x2E, 0x8B, 0xE6), "GOLD": RGBColor(0xF5, 0xC5, 0x18),
    },
    "emerald": {
        "label": "Emerald", "description": "Dark forest green with a warm brass accent.",
        "BG": RGBColor(0x0A, 0x1A, 0x14), "PANEL": RGBColor(0x11, 0x28, 0x1F),
        "LINE": RGBColor(0x1A, 0x3A, 0x2C), "ACCENT": RGBColor(0x2E, 0xB6, 0x7A), "GOLD": RGBColor(0xD9, 0xA5, 0x2C),
    },
    "charcoal": {
        "label": "Charcoal", "description": "Neutral graphite with a crisp cyan accent.",
        "BG": RGBColor(0x14, 0x15, 0x17), "PANEL": RGBColor(0x1F, 0x21, 0x24),
        "LINE": RGBColor(0x30, 0x33, 0x38), "ACCENT": RGBColor(0x4D, 0xC9, 0xE6), "GOLD": RGBColor(0xE8, 0xB4, 0x4A),
    },
    "burgundy": {
        "label": "Burgundy", "description": "Rich wine-red with a soft rose-gold accent.",
        "BG": RGBColor(0x1A, 0x0B, 0x10), "PANEL": RGBColor(0x28, 0x13, 0x1A),
        "LINE": RGBColor(0x40, 0x1F, 0x28), "ACCENT": RGBColor(0xC4, 0x5B, 0x6E), "GOLD": RGBColor(0xE0, 0xB0, 0x77),
    },
}
DEFAULT_THEME = "navy"

# Each export request runs export_project_v2() synchronously on its own
# FastAPI threadpool worker thread, so a plain module global would leak
# between concurrently-running exports with different templates -- thread-
# local storage keeps each request's selected theme isolated.
_theme_store = threading.local()


def use_theme(theme_id: str | None) -> None:
    """Selects the color palette for the deck about to be built, for the
    current thread only. Call once at the start of export_project_v2()."""
    _theme_store.theme_id = theme_id if theme_id in THEMES else DEFAULT_THEME


def _palette() -> dict:
    return THEMES[getattr(_theme_store, "theme_id", DEFAULT_THEME)]


def __getattr__(name):
    # Lets existing call sites elsewhere keep using `deck_theme.BG` /
    # `T.ACCENT` etc. as if they were plain module constants, while actually
    # resolving to the current thread's selected template (PEP 562).
    palette = _palette()
    if name in palette:
        return palette[name]
    raise AttributeError(f"module {__name__!r} has no attribute {name!r}")


# ---- fonts ----
FONT = "Calibri"          # simple, professional, widely available
FONT_LIGHT = "Calibri Light"
MM_FONT = "Pyidaungsu"    # Myanmar

BURMESE_RE = re.compile(r"[က-႟ꩠ-ꩿꧠ-꧿]")

# 16:9 widescreen (13.33 x 7.5 in)
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)


def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank(prs):
    # layout 6 = Blank
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    return prs.slides.add_slide(layout)


def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _no_autosize(tf):
    tf.word_wrap = True
    try:
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass


def add_slide(prs, page=None, company_name="ZARNI AUNG & SONS Co.,Ltd"):
    """Create a themed dark slide with pattern + footer; return the slide."""
    slide = _blank(prs)
    # dark background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _palette()["BG"]

    _add_pattern(slide, prs)
    if page is not None:
        _add_footer(slide, prs, page, company_name)
    return slide


def _add_pattern(slide, prs):
    """Subtle technical pattern: faint concentric rings + thin diagonal lines."""
    pal = _palette()
    W, H = prs.slide_width, prs.slide_height
    # large faint ring bottom-right
    for r in (Inches(2.6), Inches(3.6)):
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, W - r, H - r, r * 2, r * 2)
        c.fill.background()
        c.line.color.rgb = pal["LINE"]
        c.line.width = Pt(1)
    # faint thin lines top-left
    for i in range(3):
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(-0.5), Inches(0.4 + i * 0.28), Inches(2.2), Pt(1.2))
        ln.rotation = 35
        _solid(ln, pal["LINE"])
    # small accent dots
    for i in range(3):
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.55 + i * 0.22), Inches(0.55), Pt(5), Pt(5))
        _solid(d, pal["ACCENT"] if i == 0 else pal["LINE"])


def _add_footer(slide, prs, page, company_name):
    W, H = prs.slide_width, prs.slide_height
    # accent baseline
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), H - Inches(0.5), W, Pt(2))
    _solid(bar, _palette()["ACCENT"])
    add_text(slide, company_name, Inches(0.5), H - Inches(0.45), Inches(7), Inches(0.35),
             size=9, color=MUTED, italic=True)
    add_text(slide, str(page), W - Inches(1.2), H - Inches(0.45), Inches(0.7), Inches(0.35),
             size=9, color=MUTED, align=PP_ALIGN.RIGHT, bold=True)


def _apply_run(run, text, *, size, bold, italic, color, font):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = MM_FONT if BURMESE_RE.search(text) else font


def add_text(slide, text, left, top, width, height, *, size=16, bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT, font=FONT, anchor=MSO_ANCHOR.TOP, line_spacing=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    _no_autosize(tf)
    tf.vertical_anchor = anchor
    # support explicit newlines as separate paragraphs
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        _apply_run(run, line, size=size, bold=bold, italic=italic, color=color, font=font)
    return box


def add_title(slide, prs, title, subtitle=None):
    """Standard slide title block (top-left) with an accent tick."""
    tick = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.55), Pt(6), Inches(0.7))
    _solid(tick, _palette()["GOLD"])
    add_text(slide, title, Inches(0.75), Inches(0.45), prs.slide_width - Inches(1.3), Inches(0.9),
             size=28, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.78), Inches(1.28), prs.slide_width - Inches(1.3), Inches(0.5),
                 size=15, italic=True, color=_palette()["ACCENT"])
        return Inches(1.95)
    return Inches(1.6)


def add_richtext(slide, blocks, left, top, width, height, *, size=15, color=WHITE):
    """Render parsed rich-text blocks (see richtext.parse_html) into a text frame.
    blocks: list of {type: 'p'|'ul'|'ol', items: [ [ {text,bold,italic}, ... ], ... ] }"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    _no_autosize(tf)
    first = True
    for block in blocks:
        btype = block.get("type", "p")
        for idx, runs in enumerate(block.get("items", [])):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(6)
            p.line_spacing = 1.1
            prefix = ""
            if btype == "ul":
                prefix = "•  "
            elif btype == "ol":
                prefix = f"{idx + 1}.  "
            if prefix:
                pr = p.add_run()
                _apply_run(pr, prefix, size=size, bold=True, italic=False, color=_palette()["ACCENT"], font=FONT)
            for r in runs:
                run = p.add_run()
                _apply_run(run, r.get("text", ""), size=size, bold=r.get("bold", False),
                           italic=r.get("italic", False), color=color, font=FONT)
    return box


def add_panel(slide, left, top, width, height, color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    _solid(shp, color if color is not None else _palette()["PANEL"])
    return shp


def add_image_contain(slide, image_stream, left, top, box_w, box_h):
    """Place an image preserving aspect ratio within a box, centered."""
    from PIL import Image
    import io
    data = image_stream if isinstance(image_stream, (bytes, bytearray)) else image_stream.read()
    try:
        with Image.open(io.BytesIO(data)) as im:
            iw, ih = im.size
    except Exception:
        iw, ih = int(box_w), int(box_h)
    scale = min(box_w / iw, box_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    import io as _io
    slide.shapes.add_picture(_io.BytesIO(data), int(left + (box_w - w) / 2), int(top + (box_h - h) / 2), w, h)
