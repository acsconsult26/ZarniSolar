"""Redesigned proposal deck (v2): dark-themed, built from scratch from form
data. Currently builds slides 1-12; later sections will be added incrementally.
"""
from __future__ import annotations

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from ..schema import merged_field_values
from . import deck_theme as T
from .richtext import parse_html

CONTENTS_ITEMS = [
    "Project Objectives",
    "Surveying Data (Project Background)",
    "System Requirement",
    "Technical Proposal",
    "Product Specifications",
    "Technical Advantages",
    "Solar Panels Support Mounting Structure",
    "Warranty",
]


def _fmt(v, suffix=""):
    if v is None or v == "":
        return "—"
    if isinstance(v, float) and v.is_integer():
        v = int(v)
    return f"{v}{suffix}"


def _money(v):
    if v is None or v == "":
        return "—"
    try:
        return f"{float(v):,.0f}"
    except (TypeError, ValueError):
        return str(v)


class _Imgs:
    def __init__(self, uploads, storage):
        self.uploads = uploads or {}
        self.storage = storage

    def get(self, field):
        path = self.uploads.get(field)
        if path and self.storage.exists(path):
            return self.storage.read_bytes(path)
        return None


def export_project_v2(project, storage, company_info=None, selected_products=None,
                      closing_statement=None, warranty_lines=None) -> bytes:
    data = project.data or {}
    v = merged_field_values(data)
    ci = company_info or {}
    company_name = ci.get("company_name") or "ZARNI AUNG & SONS Co.,Ltd"
    contact = ci.get("contact") or "09-2031977"
    client = v.get("site_name") or "Client"
    imgs = _Imgs(project.uploads, storage)

    prs = T.new_deck()
    W = prs.slide_width

    _slide1_cover(prs, client, v.get("proposal_date"), company_name, contact)
    _slide2_intro(prs, v.get("introduction"), company_name)
    _slide3_contents(prs, company_name)
    _slide4_objectives(prs, v.get("project_objectives"), company_name)
    _slide5_surveying(prs, v, imgs, company_name)
    _slide6_bill(prs, client, v, imgs, company_name)
    _slide7_survey_photos(prs, imgs, company_name)
    _slide8_solar_frame(prs, imgs, company_name)
    _slide9_data_result(prs, client, v, imgs, company_name, page=9,
                        prefix="", title=f"Surveying Data Result for {client}")
    _slide10_analyzer(prs, client, v.get("analyzer_date_range"), imgs.get("analyzer_image"),
                      company_name, page=10)

    if imgs.get("analyzer_chart"):
        page = len(prs.slides._sldIdLst) + 1
        _slide_analyzer_stats(prs, f"Power Analyzer — Hourly Load Profile for {client}",
                              data.get("analyzer_stats"), imgs.get("analyzer_chart"), company_name, page)

    if data.get("include_second_survey"):
        _slide9_data_result(prs, client, v, imgs, company_name, page=11, prefix="second_",
                            title="Surveying Data Result (Other Meters)")

    # Slides 13-14 : System Requirement options (up to 4, 2 per slide)
    options = [o for o in (data.get("system_options") or [])
               if (o.get("items") or o.get("capex"))]
    if options:
        page = len(prs.slides._sldIdLst) + 1
        _slide_options(prs, "System Requirement", options[:2], page, company_name)
    if len(options) > 2:
        page = len(prs.slides._sldIdLst) + 1
        _slide_options(prs, "System Requirement (cont.)", options[2:4], page, company_name)

    # Slide 15 : ROI (EPC-only vs Solar) - only if the ROI inputs are present
    if data.get("roi_total_epc_units") and data.get("roi_avg_unit_cost"):
        page = len(prs.slides._sldIdLst) + 1
        _slide_roi(prs, v, company_name, page)

    # Slide 16 : daily power usage comparison chart (Load / Grid / Solar per option)
    chart_opts = [o for o in options if (o.get("grid_units") not in (None, "")
                                         or o.get("solar_units") not in (None, ""))]
    if chart_opts:
        page = len(prs.slides._sldIdLst) + 1
        _slide_usage_chart(prs, data, chart_opts, company_name, page)

    # Slide 17 : payback comparison table
    if options and (data.get("payback_epc_units_month") or any(o.get("payback_years") for o in options)):
        page = len(prs.slides._sldIdLst) + 1
        _slide_payback(prs, data, options[:4], company_name, page)

    # Slide 18 : Hybrid Solar System Drawing (image upload)
    if imgs.get("system_drawing_image"):
        page = len(prs.slides._sldIdLst) + 1
        _slide_single_image(prs, "Hybrid Solar System Drawing", imgs.get("system_drawing_image"),
                            company_name, page)

    # Slide 19 : Solar System Block Diagram (image upload)
    if imgs.get("block_diagram_image"):
        page = len(prs.slides._sldIdLst) + 1
        _slide_single_image(prs, "Solar System Block Diagram", imgs.get("block_diagram_image"),
                            company_name, page)

    # Slide 20 : Simulation Result (South View, West View, East View, PV Array)
    sim_imgs = [imgs.get("sim_south_view"), imgs.get("sim_west_view"), imgs.get("sim_east_view"), imgs.get("sim_pv_array")]
    if any(sim_imgs):
        page = len(prs.slides._sldIdLst) + 1
        _slide_photo_row(prs, "Simulation Result", sim_imgs, company_name, page,
                         labels=["South View", "West View", "East View", "PV Array"])

    # Slide 21 : Energy Yield Report (2 photos side by side)
    eyr_imgs = [imgs.get("energy_yield_1"), imgs.get("energy_yield_2")]
    if any(eyr_imgs):
        page = len(prs.slides._sldIdLst) + 1
        _slide_photo_row(prs, "Energy Yield Report", eyr_imgs, company_name, page)

    # Slide 22 : Monthly Production From Solar (single full-size image)
    if imgs.get("monthly_production_image"):
        page = len(prs.slides._sldIdLst) + 1
        _slide_single_image(prs, "Monthly Production From Solar", imgs.get("monthly_production_image"),
                            company_name, page)

    # Slide 23 : West View Shade Report (Perfect?) - 2 photos
    west_imgs = [imgs.get("west_shade_1"), imgs.get("west_shade_2")]
    if any(west_imgs):
        page = len(prs.slides._sldIdLst) + 1
        title = "West View Shade Report" + (" (Perfect)" if data.get("west_shade_perfect") else "")
        _slide_photo_row(prs, title, west_imgs, company_name, page)

    # Slide 24 : South View Shade Report (Perfect?) - 2 photos
    south_imgs = [imgs.get("south_shade_1"), imgs.get("south_shade_2")]
    if any(south_imgs):
        page = len(prs.slides._sldIdLst) + 1
        title = "South View Shade Report" + (" (Perfect)" if data.get("south_shade_perfect") else "")
        _slide_photo_row(prs, title, south_imgs, company_name, page)

    # Slide 24b : East View Shade Report (Perfect?) - 2 photos
    east_imgs = [imgs.get("east_shade_1"), imgs.get("east_shade_2")]
    if any(east_imgs):
        page = len(prs.slides._sldIdLst) + 1
        title = "East View Shade Report" + (" (Perfect)" if data.get("east_shade_perfect") else "")
        _slide_photo_row(prs, title, east_imgs, company_name, page)

    # Slide 25 : fixed FYI info about Zarni Electronics Service (admin-editable)
    if closing_statement:
        page = len(prs.slides._sldIdLst) + 1
        _slide_closing_statement(prs, closing_statement, company_name, page)

    # Slides 26-28 : selected product specifications (Solar, Battery, Inverter)
    for category, title in (("panel", "Solar Panel Specification"),
                            ("battery", "Battery Specification"),
                            ("inverter", "Inverter Specification"),
                            ("gateway", "Gateway Specification")):
        product = (selected_products or {}).get(category)
        if product:
            page = len(prs.slides._sldIdLst) + 1
            image = _product_image(storage, product)
            _slide_product_spec(prs, title, product, image, company_name, page)

    # Slides 29-30 : Technical Advantages (free-form, filled per-proposal)
    for key in ("tech_advantages_1", "tech_advantages_2"):
        html = data.get(key)
        if html:
            page = len(prs.slides._sldIdLst) + 1
            _slide_richtext_block(prs, "Technical Advantages", html, company_name, page)

    # Slide 31 : Solar Support Mounting Structure -- design images
    mount_design_imgs = [imgs.get("mounting_design_1"), imgs.get("mounting_design_2")]
    if any(mount_design_imgs):
        page = len(prs.slides._sldIdLst) + 1
        _slide_photo_row(prs, "Solar Support Mounting Structure — Design", mount_design_imgs, company_name, page)

    # Slides 32-33 : Solar Priority Install Area 1 & 2 (image + note)
    for idx, label in ((1, "Solar Priority Install Area 1"), (2, "Solar Priority Install Area 2")):
        img = imgs.get(f"mounting_area_{idx}_image")
        note = data.get(f"mounting_area_{idx}_note")
        if img or note:
            page = len(prs.slides._sldIdLst) + 1
            _slide_image_note(prs, label, img, note, company_name, page)

    # Slide 34 : Warranty (admin-managed template, chosen per-proposal)
    if warranty_lines:
        page = len(prs.slides._sldIdLst) + 1
        _slide_bullets(prs, "Warranty", warranty_lines, company_name, page)

    import io
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


# ---------------- individual slides ----------------

def _slide1_cover(prs, client, date, company_name, contact):
    slide = T.add_slide(prs, page=None, company_name=company_name)
    W, H = prs.slide_width, prs.slide_height
    # accent block
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(2.35), Inches(0.12), Inches(2.0))
    T._solid(bar, T.GOLD)
    T.add_text(slide, "BUSINESS ENERGY SOLUTION", Inches(1.15), Inches(2.3), W - Inches(2), Inches(0.7),
               size=34, bold=True, color=T.WHITE)
    T.add_text(slide, f"For  {client}", Inches(1.15), Inches(3.05), W - Inches(2), Inches(0.8),
               size=30, bold=True, color=T.ACCENT)
    T.add_text(slide, "Technical Proposal", Inches(1.18), Inches(3.95), W - Inches(2), Inches(0.5),
               size=18, italic=True, color=T.MUTED)
    T.add_text(slide, company_name, Inches(1.18), Inches(4.4), W - Inches(2), Inches(0.5),
               size=16, bold=True, color=T.WHITE)
    # date + contact bottom
    T.add_text(slide, f"Date – {date or ''}", Inches(0.9), H - Inches(1.1), Inches(5), Inches(0.4),
               size=13, color=T.MUTED)
    T.add_text(slide, f"Contact To; {contact}", W - Inches(5.5), H - Inches(1.1), Inches(4.6), Inches(0.4),
               size=13, color=T.MUTED, align=PP_ALIGN.RIGHT)


def _slide2_intro(prs, introduction, company_name):
    slide = T.add_slide(prs, page=2, company_name=company_name)
    top = T.add_title(slide, prs, "Introduction")
    blocks = parse_html(introduction) or [{"type": "p", "items": [[{"text": "", "bold": False, "italic": False}]]}]
    T.add_richtext(slide, blocks, Inches(0.9), top, prs.slide_width - Inches(1.8),
                   prs.slide_height - top - Inches(0.8), size=16)


def _slide3_contents(prs, company_name):
    slide = T.add_slide(prs, page=3, company_name=company_name)
    top = T.add_title(slide, prs, "CONTENTS")
    box = slide.shapes.add_textbox(Inches(1.2), top, prs.slide_width - Inches(2.4), prs.slide_height - top - Inches(0.8))
    tf = box.text_frame
    T._no_autosize(tf)
    for i, item in enumerate(CONTENTS_ITEMS):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        r1 = p.add_run(); T._apply_run(r1, f"{i + 1}.  ", size=20, bold=True, italic=False, color=T.GOLD, font=T.FONT)
        r2 = p.add_run(); T._apply_run(r2, item, size=20, bold=False, italic=False, color=T.WHITE, font=T.FONT)


def _slide4_objectives(prs, objectives, company_name):
    slide = T.add_slide(prs, page=4, company_name=company_name)
    top = T.add_title(slide, prs, "Project Objectives  –  ရည်ရွယ်ချက်များ")
    blocks = parse_html(objectives) or [{"type": "p", "items": [[{"text": "", "bold": False, "italic": False}]]}]
    T.add_richtext(slide, blocks, Inches(0.9), top, prs.slide_width - Inches(1.8),
                   prs.slide_height - top - Inches(0.8), size=15)


def _slide5_surveying(prs, v, imgs, company_name):
    slide = T.add_slide(prs, page=5, company_name=company_name)
    top = T.add_title(slide, prs, "Surveying Data  –  Project Background")
    lat, lng = v.get("survey_lat"), v.get("survey_lng")
    loc = f"{lat}, {lng}" if (lat or lng) else "—"
    rows = [
        ("Project Location", loc),
        ("Project Solution", v.get("project_solution") or "—"),
        ("Solar Panel Installation Area", _fmt(v.get("install_area_sqft"), " sq ft")),
        ("Tilt Angle", _fmt(v.get("tilt_angle"), "°")),
    ]
    box = slide.shapes.add_textbox(Inches(0.9), top, Inches(6.2), prs.slide_height - top - Inches(0.9))
    tf = box.text_frame; T._no_autosize(tf)
    for i, (k, val) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        r1 = p.add_run(); T._apply_run(r1, f"{k}:  ", size=16, bold=True, italic=False, color=T.ACCENT, font=T.FONT)
        r2 = p.add_run(); T._apply_run(r2, str(val), size=16, bold=False, italic=False, color=T.WHITE, font=T.FONT)
    img = imgs.get("survey_image")
    if img:
        T.add_image_contain(slide, img, Inches(7.4), top, Inches(5.0), prs.slide_height - top - Inches(1.0))


def _slide6_bill(prs, client, v, imgs, company_name):
    slide = T.add_slide(prs, page=6, company_name=company_name)
    top = T.add_title(slide, prs, f"{client}  ဓာတ်အားခ သုံးစွဲယူနစ်တောင်းခံလွှာ")
    img = imgs.get("meter_bill_image")
    if img:
        T.add_image_contain(slide, img, Inches(0.9), top, Inches(5.6), prs.slide_height - top - Inches(1.0))
    # figures panel on the right
    px = Inches(7.0)
    rows = [
        ("ဓာတ်အားခ စုစုပေါင်း", f"{_money(v.get('total_epc_cost'))} ကျပ်"),
        ("သုံးစွဲယူနစ် စုစုပေါင်း", f"{_money(v.get('total_epc_units'))} ယူနစ်"),
        ("တစ်ယူနစ် ကျသင့်ငွေ", f"{_money(v.get('per_unit_cost'))} ကျပ်"),
    ]
    box = slide.shapes.add_textbox(px, top + Inches(0.3), Inches(5.3), Inches(4))
    tf = box.text_frame; T._no_autosize(tf)
    for i, (k, val) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(18)
        r1 = p.add_run(); T._apply_run(r1, f"{k}  –  ", size=17, bold=False, italic=False, color=T.MUTED, font=T.MM_FONT)
        r2 = p.add_run(); T._apply_run(r2, str(val), size=18, bold=True, italic=False, color=T.GOLD, font=T.MM_FONT)


def _slide7_survey_photos(prs, imgs, company_name):
    slide = T.add_slide(prs, page=7, company_name=company_name)
    top = T.add_title(slide, prs, "Surveying ပြုလုပ်သည့် မှတ်တမ်းပုံများ")
    half = (prs.slide_width - Inches(2.1)) / 2
    for i, field in enumerate(("survey_photo_1", "survey_photo_2")):
        left = Inches(0.7) + i * (half + Inches(0.5))
        img = imgs.get(field)
        if img:
            T.add_image_contain(slide, img, left, top, half, prs.slide_height - top - Inches(1.0))
        else:
            _placeholder(slide, left, top, half, prs.slide_height - top - Inches(1.0))


def _slide8_solar_frame(prs, imgs, company_name):
    slide = T.add_slide(prs, page=8, company_name=company_name)
    top = T.add_title(slide, prs, "Solar Frame ပြုလုပ်မည့် နေရာ")
    fields = [("solar_frame_map_image", "Google Map Overview"),
              ("solar_frame_site_1", "Site View 1"),
              ("solar_frame_site_2", "Site View 2")]
    third = (prs.slide_width - Inches(2.4)) / 3
    h = prs.slide_height - top - Inches(1.2)
    for i, (field, label) in enumerate(fields):
        left = Inches(0.7) + i * (third + Inches(0.45))
        img = imgs.get(field)
        if img:
            T.add_image_contain(slide, img, left, top, third, h)
        else:
            _placeholder(slide, left, top, third, h)
        T.add_text(slide, label, left, top + h + Inches(0.1), third, Inches(0.4),
                   size=12, italic=True, color=T.MUTED, align=PP_ALIGN.CENTER)


def _slide9_data_result(prs, client, v, imgs, company_name, page, prefix, title):
    slide = T.add_slide(prs, page=page, company_name=company_name)
    top = T.add_title(slide, prs, title)

    def g(base):
        return v.get(f"{prefix}{base}") if prefix else v.get(base)

    rows = []
    if prefix:
        # Only the optional Second Survey still collects these manually --
        # the primary survey now derives consumption from the uploaded
        # Excel/CSV file instead of asking for them.
        rows.append(("Maximum Load Consumption", _fmt(g("max_load_kw"), " kWp")))
        rows.append(("Duration Hours", _fmt(g("duration_hours"), " Hr")))
    rows.append(("Line Voltage", _fmt(g("voltage_v"), " V")))
    if prefix:
        rows.append(("Power Factor", _fmt(g("power_factor"))))
    rows.append(("Transformer Size", _fmt(g("transformer_kva"), " kVA")))
    rows.append(("Generator", _fmt(g("generator_capacity_kva") if not prefix else g("generator_kva"), " kVA")))
    rows.append(("PV Installation Area", _fmt(g("pv_installation_area_sqft") if not prefix else g("pv_area_sqft"), " sq ft")))
    rows.append(("Average Consumption", _fmt(g("avg_units") if prefix else v.get("survey_avg_units"), " Units/day")))
    rows.append(("Peak Consumption", _fmt(g("peak_units") if prefix else v.get("survey_peak_units"), " Units")))
    box = slide.shapes.add_textbox(Inches(0.9), top, Inches(7.2), prs.slide_height - top - Inches(0.9))
    tf = box.text_frame; T._no_autosize(tf)
    for i, (k, val) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        r1 = p.add_run(); T._apply_run(r1, f"{k}  –  ", size=15, bold=True, italic=False, color=T.ACCENT, font=T.FONT)
        r2 = p.add_run(); T._apply_run(r2, str(val), size=15, bold=False, italic=False, color=T.WHITE, font=T.FONT)


def _slide10_analyzer(prs, client, date_range, image, company_name, page):
    slide = T.add_slide(prs, page=page, company_name=company_name)
    top = T.add_title(slide, prs, f"Power Analyzer Recording for {client}",
                      subtitle=(date_range or None))
    if image:
        T.add_image_contain(slide, image, Inches(0.9), top, prs.slide_width - Inches(1.8),
                            prs.slide_height - top - Inches(0.9))
    else:
        _placeholder(slide, Inches(0.9), top, prs.slide_width - Inches(1.8), prs.slide_height - top - Inches(0.9))


def _slide_analyzer_stats(prs, title, stats, chart_image, company_name, page):
    slide = T.add_slide(prs, page=page, company_name=company_name)
    top = T.add_title(slide, prs, title)

    stats = stats or {}

    def stat_line(label, avg_key, peak_key, unit):
        avg = stats.get(avg_key)
        peak = stats.get(peak_key)
        avg_s = f"{avg}{unit}" if avg is not None else "—"
        peak_s = f"{peak}{unit}" if peak is not None else "—"
        return f"{label}   Avg {avg_s}   ·   Peak {peak_s}"

    lines = [
        stat_line("kW", "avg_kw", "peak_kw", " kW"),
        stat_line("PF", "avg_pf", "peak_pf", ""),
        stat_line("THD (Voltage)", "avg_thd_voltage", "peak_thd_voltage", "%"),
        stat_line("THD (Current)", "avg_thd_current", "peak_thd_current", "%"),
    ]
    T.add_text(slide, "\n".join(lines), Inches(0.6), top, prs.slide_width - Inches(1.2), Inches(1.3),
              size=14, color=T.MUTED, line_spacing=1.3)

    chart_top = top + Inches(1.4)
    if chart_image:
        T.add_image_contain(slide, chart_image, Inches(0.6), chart_top, prs.slide_width - Inches(1.2),
                            prs.slide_height - chart_top - Inches(0.6))
    else:
        _placeholder(slide, Inches(0.6), chart_top, prs.slide_width - Inches(1.2),
                     prs.slide_height - chart_top - Inches(0.6))


def _slide_photo_row(prs, title, images, company_name, page, labels=None):
    """Render 2-3 images side by side, equal width, with optional captions."""
    slide = T.add_slide(prs, page=page, company_name=company_name)
    top = T.add_title(slide, prs, title)
    n = len(images)
    gap = Inches(0.45)
    margin = Inches(0.7)
    w = (prs.slide_width - margin * 2 - gap * (n - 1)) / n
    h = prs.slide_height - top - (Inches(1.2) if labels else Inches(1.0))
    for i, img in enumerate(images):
        left = margin + i * (w + gap)
        if img:
            T.add_image_contain(slide, img, left, top, w, h)
        else:
            _placeholder(slide, left, top, w, h)
        if labels:
            T.add_text(slide, labels[i], left, top + h + Inches(0.1), w, Inches(0.4),
                       size=12, italic=True, color=T.MUTED, align=PP_ALIGN.CENTER)


def _slide_closing_statement(prs, text, company_name, page):
    slide = T.add_slide(prs, page=page, company_name=company_name)
    top = T.add_title(slide, prs, "Zarni Electronics Service")
    blocks = parse_html(text)
    T.add_richtext(slide, blocks, Inches(0.9), top, prs.slide_width - Inches(1.8),
                   prs.slide_height - top - Inches(0.8), size=16)


def _slide_richtext_block(prs, title, html, company_name, page):
    """Generic single-richtext-block slide (Technical Advantages 29-30)."""
    slide = T.add_slide(prs, page=page, company_name=company_name)
    top = T.add_title(slide, prs, title)
    blocks = parse_html(html)
    T.add_richtext(slide, blocks, Inches(0.9), top, prs.slide_width - Inches(1.8),
                   prs.slide_height - top - Inches(0.8), size=16)


def _slide_image_note(prs, title, image, note, company_name, page):
    """Image (left) + a short caption/measurement note (right) -- the two
    Solar Priority Install Area slides (32-33)."""
    slide = T.add_slide(prs, page=page, company_name=company_name)
    top = T.add_title(slide, prs, title)
    body_h = prs.slide_height - top - Inches(0.8)
    if image:
        T.add_image_contain(slide, image, Inches(0.7), top, Inches(7.5), body_h)
    else:
        _placeholder(slide, Inches(0.7), top, Inches(7.5), body_h)
    if note:
        T.add_text(slide, note, Inches(8.5), top, prs.slide_width - Inches(9.2), body_h,
                   size=18, bold=True, color=T.WHITE, anchor=MSO_ANCHOR.MIDDLE)


def _slide_bullets(prs, title, lines, company_name, page):
    """Simple bulleted-line slide (Warranty, slide 34)."""
    slide = T.add_slide(prs, page=page, company_name=company_name)
    top = T.add_title(slide, prs, title)
    box = slide.shapes.add_textbox(Inches(0.9), top, prs.slide_width - Inches(1.8),
                                   prs.slide_height - top - Inches(0.9))
    tf = box.text_frame
    T._no_autosize(tf)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        run = p.add_run()
        T._apply_run(run, f"•  {line}", size=17, bold=False, italic=False, color=T.WHITE, font=T.FONT)


def _product_image(storage, product):
    path = (product or {}).get("image_path")
    if path and storage.exists(path):
        return storage.read_bytes(path)
    return None


def _slide_product_spec(prs, title, product, image, company_name, page):
    slide = T.add_slide(prs, page=page, company_name=company_name)
    subtitle = product.get("spec_title") or f"{product.get('brand', '')} {product.get('model_name', '')}".strip()
    top = T.add_title(slide, prs, title, subtitle=subtitle or None)

    img_w = Inches(4.2)
    tbl_left = Inches(0.6) + img_w + Inches(0.4)
    tbl_width = prs.slide_width - tbl_left - Inches(0.6)
    height = prs.slide_height - top - Inches(0.8)

    if image:
        T.add_image_contain(slide, image, Inches(0.6), top, img_w, height)
    else:
        _placeholder(slide, Inches(0.6), top, img_w, height)

    specs = product.get("specs") or []
    nrows = max(len(specs), 1)
    gf = slide.shapes.add_table(nrows, 2, int(tbl_left), int(top), int(tbl_width), int(height))
    table = gf.table
    table.first_row = False
    table.horz_banding = False
    table.columns[0].width = int(tbl_width * 0.55)
    table.columns[1].width = int(tbl_width * 0.45)

    def cell(r, c, text, *, bold, color, bg):
        cl = table.cell(r, c)
        cl.fill.solid(); cl.fill.fore_color.rgb = bg
        cl.vertical_anchor = MSO_ANCHOR.MIDDLE
        cl.margin_left = Inches(0.1)
        tf = cl.text_frame; tf.word_wrap = True
        run = tf.paragraphs[0].add_run()
        T._apply_run(run, text, size=13, bold=bold, italic=False, color=color, font=T.FONT)

    if specs:
        for i, s in enumerate(specs):
            bg = T.PANEL if i % 2 == 0 else T.BG
            cell(i, 0, s.get("label", ""), bold=True, color=T.ACCENT, bg=bg)
            cell(i, 1, f"{s.get('value', '')} {s.get('unit', '')}".strip(), bold=False, color=T.WHITE, bg=bg)
    else:
        cell(0, 0, "No specifications added", bold=False, color=T.MUTED, bg=T.PANEL)
        cell(0, 1, "", bold=False, color=T.WHITE, bg=T.PANEL)


def _slide_single_image(prs, title, image, company_name, page):
    slide = T.add_slide(prs, page=page, company_name=company_name)
    top = T.add_title(slide, prs, title)
    if image:
        T.add_image_contain(slide, image, Inches(0.6), top, prs.slide_width - Inches(1.2),
                            prs.slide_height - top - Inches(0.8))
    else:
        _placeholder(slide, Inches(0.6), top, prs.slide_width - Inches(1.2), prs.slide_height - top - Inches(0.8))


def _option_line(item):
    name = (item.get("name") or "").strip()
    qty = item.get("qty")
    unit = (item.get("unit") or "Nos").strip()
    if qty in (None, ""):
        return name
    return f"{name}  ({qty} {unit})"


def _slide_options(prs, title, options, page, company_name):
    slide = T.add_slide(prs, page=page, company_name=company_name)
    top = T.add_title(slide, prs, title)
    ncols = len(options)
    max_items = max((len(o.get("items") or []) for o in options), default=0)
    nrows = 1 + max_items + 1  # title header + item rows + CAPEX row

    left = Inches(0.7)
    width = prs.slide_width - Inches(1.4)
    height = prs.slide_height - top - Inches(0.7)
    gf = slide.shapes.add_table(nrows, ncols, left, top, int(width), int(height))
    table = gf.table
    table.first_row = False
    table.horz_banding = False
    col_w = int(width / ncols)
    for c in range(ncols):
        table.columns[c].width = col_w

    def cell(r, c, text, *, size, bold, color, bg):
        cl = table.cell(r, c)
        cl.fill.solid(); cl.fill.fore_color.rgb = bg
        cl.vertical_anchor = MSO_ANCHOR.MIDDLE
        cl.margin_top = Pt(3); cl.margin_bottom = Pt(3)
        tf = cl.text_frame; tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        T._apply_run(run, text, size=size, bold=bold, italic=False, color=color, font=T.FONT)

    for ci, opt in enumerate(options):
        cell(0, ci, opt.get("title") or f"Option {ci + 1}", size=16, bold=True, color=T.WHITE, bg=T.ACCENT)
        items = opt.get("items") or []
        for ri in range(max_items):
            txt = _option_line(items[ri]) if ri < len(items) else ""
            cell(ri + 1, ci, txt, size=12, bold=False, color=T.WHITE, bg=T.PANEL)
        capex = opt.get("capex")
        capex_txt = f"Est. CAPEX = {_money(capex)} MMK" if capex not in (None, "") else "Est. CAPEX = —"
        cell(nrows - 1, ci, capex_txt, size=14, bold=True, color=T.BG, bg=T.GOLD)


def _num(v):
    try:
        return float(v)
    except (TypeError, ValueError):
        return 0.0


def _slide_roi(prs, v, company_name, page):
    import math

    n1 = _num(v.get("roi_total_epc_units"))
    n2 = _num(v.get("roi_epc_with_solar_units"))
    n3 = _num(v.get("roi_solar_units"))
    cost = _num(v.get("roi_avg_unit_cost"))
    years = int(_num(v.get("roi_years")) or 5)

    annual_epc = n1 * 30 * 12 * cost
    annual_solar = n3 * 30 * 12 * cost
    total_epc = annual_epc * years
    total_solar = annual_solar * years
    savings = total_epc - total_solar

    def i(x):
        return f"{x:,.0f}"

    slide = T.add_slide(prs, page=page, company_name=company_name)
    top = T.add_title(slide, prs, "ROI  –  ရင်းနှီးမြှုပ်နှံမှု ပြန်လည်ရရှိမှု တွက်ချက်မှု")

    # summary / ratio line
    ratio_txt = "—"
    if n2 and n3:
        g = math.gcd(int(n2), int(n3)) or 1
        epc_pct = (n2 / n1 * 100) if n1 else 0
        solar_pct = (n3 / n1 * 100) if n1 else 0
        ratio_txt = (f"စုစုပေါင်း သုံးစွဲမှု = {i(n1)} ယူနစ်     |     "
                     f"EPC : Solar = {int(n2 // g)} : {int(n3 // g)}  "
                     f"({epc_pct:.0f}% : {solar_pct:.0f}%)")
    T.add_text(slide, ratio_txt, Inches(0.9), top, prs.slide_width - Inches(1.8), Inches(0.5),
               size=15, color=T.GOLD, italic=True)
    tbl_top = top + Inches(0.65)

    # 3 rows x 2 cols table
    left = Inches(0.7)
    width = prs.slide_width - Inches(1.4)
    height = prs.slide_height - tbl_top - Inches(1.2)
    gf = slide.shapes.add_table(3, 2, left, tbl_top, int(width), int(height))
    table = gf.table
    table.first_row = False
    table.horz_banding = False
    table.columns[0].width = int(width / 2)
    table.columns[1].width = int(width / 2)

    def mcell(r, c, lines, *, size, bold, color, bg):
        cl = table.cell(r, c)
        cl.fill.solid(); cl.fill.fore_color.rgb = bg
        cl.vertical_anchor = MSO_ANCHOR.MIDDLE
        cl.margin_left = Inches(0.15); cl.margin_right = Inches(0.15)
        tf = cl.text_frame; tf.word_wrap = True
        for k, ln in enumerate(lines):
            p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            T._apply_run(run, ln, size=size, bold=bold, italic=False, color=color, font=T.FONT)

    mcell(0, 0, ["EPC ဖြင့်သာ သုံးစွဲပါက"], size=16, bold=True, color=T.WHITE, bg=T.ACCENT)
    mcell(0, 1, ["Solar ဖြင့် သုံးစွဲပါက"], size=16, bold=True, color=T.WHITE, bg=T.ACCENT)

    mcell(1, 0, [
        "တစ်နှစ်စာ ကုန်ကျငွေ",
        f"{i(n1)} ယူနစ် × ၃၀ ရက် × ၁၂ လ × {i(cost)} ကျပ်",
        f"= {i(annual_epc)} ကျပ်",
    ], size=13, bold=False, color=T.WHITE, bg=T.PANEL)
    mcell(1, 1, [
        "တစ်နှစ်စာ ကုန်ကျငွေ",
        f"{i(n3)} ယူနစ် × ၃၀ ရက် × ၁၂ လ × {i(cost)} ကျပ်",
        f"= {i(annual_solar)} ကျပ်",
    ], size=13, bold=False, color=T.WHITE, bg=T.PANEL)

    mcell(2, 0, [
        f"{years} နှစ်စာ ကုန်ကျငွေ",
        f"{i(annual_epc)} × {years}",
        f"= {i(total_epc)} ကျပ်",
    ], size=13, bold=True, color=T.GOLD, bg=T.PANEL)
    mcell(2, 1, [
        f"{years} နှစ်စာ ကုန်ကျငွေ",
        f"{i(annual_solar)} × {years}",
        f"= {i(total_solar)} ကျပ်",
    ], size=13, bold=True, color=T.GOLD, bg=T.PANEL)

    # savings conclusion
    T.add_text(slide,
               f"{years} နှစ်အတွင်း ချွေတာနိုင်မှု = {i(total_epc)} − {i(total_solar)} = {i(savings)} ကျပ်",
               Inches(0.7), prs.slide_height - Inches(1.05), width, Inches(0.5),
               size=15, bold=True, color=T.WHITE, align=PP_ALIGN.CENTER)


def _slide_payback(prs, data, options, company_name, page):
    def n(x):
        try:
            return float(x)
        except (TypeError, ValueError):
            return 0.0

    def i0(x):
        return "—" if x in (None, "") else (f"{int(float(x))}" if float(x).is_integer() else f"{float(x):g}")

    epc_month = n(data.get("payback_epc_units_month"))
    cost = n(data.get("payback_unit_cost"))
    epc_year_cost = epc_month * 12 * cost

    slide = T.add_slide(prs, page=page, company_name=company_name)
    top = T.add_title(slide, prs, "အရင်းကြေကာလ ရှင်းလင်းတင်ပြခြင်း")

    labels = ["", "System", "CAPEX (MMK)", "EPC / Month", "EPC / Year (MMK)",
              "Solar / Month", "Solar / Year", "Payback Period"]
    ncols = 1 + len(options)
    nrows = 8
    left = Inches(0.5)
    width = prs.slide_width - Inches(1.0)
    height = prs.slide_height - top - Inches(0.7)
    gf = slide.shapes.add_table(nrows, ncols, left, top, int(width), int(height))
    table = gf.table
    table.first_row = False
    table.horz_banding = False
    label_w = int(width * 0.20)
    opt_w = int((width - label_w) / len(options))
    table.columns[0].width = label_w
    for c in range(1, ncols):
        table.columns[c].width = opt_w

    def cell(r, c, text, *, size, bold, color, bg):
        cl = table.cell(r, c)
        cl.fill.solid(); cl.fill.fore_color.rgb = bg
        cl.vertical_anchor = MSO_ANCHOR.MIDDLE
        cl.margin_left = Inches(0.06); cl.margin_right = Inches(0.06)
        cl.margin_top = Pt(2); cl.margin_bottom = Pt(2)
        tf = cl.text_frame; tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        T._apply_run(run, text, size=size, bold=bold, italic=False, color=color, font=T.FONT)

    # label column
    cell(0, 0, "Options", size=12, bold=True, color=T.WHITE, bg=T.ACCENT)
    for r in range(1, nrows):
        cell(r, 0, labels[r], size=12, bold=True, color=T.WHITE, bg=T.BG)

    for ci, o in enumerate(options, start=1):
        solar_c = o.get("solar_count"); inv = o.get("inverter_power")
        batt = o.get("battery_capacity"); backup = o.get("backup_hours")
        system_txt = f"{i0(solar_c)} Solar, {i0(inv)}kW, {i0(batt)}kWh, {i0(backup)}Hrs"
        solar_units = n(o.get("solar_units"))
        solar_month = solar_units * 30
        solar_year = solar_month * 12
        py = o.get("payback_years")

        cell(0, ci, o.get("title") or f"Option {ci}", size=13, bold=True, color=T.WHITE, bg=T.ACCENT)
        cell(1, ci, system_txt, size=11, bold=False, color=T.WHITE, bg=T.PANEL)
        cell(2, ci, _money(o.get("capex")), size=12, bold=True, color=T.GOLD, bg=T.PANEL)
        cell(3, ci, f"{_money(epc_month)} Units", size=11, bold=False, color=T.WHITE, bg=T.PANEL)
        cell(4, ci, _money(epc_year_cost), size=11, bold=False, color=T.WHITE, bg=T.PANEL)
        cell(5, ci, f"{_money(solar_month)} Units", size=11, bold=False, color=T.WHITE, bg=T.PANEL)
        cell(6, ci, f"{_money(solar_year)} Units", size=11, bold=False, color=T.WHITE, bg=T.PANEL)
        cell(7, ci, f"{i0(py)} Years", size=12, bold=True, color=T.GOLD, bg=T.PANEL)


def _slide_usage_chart(prs, data, chart_opts, company_name, page):
    from .chart_usage import render_usage_chart

    def n(x):
        try:
            return float(x)
        except (TypeError, ValueError):
            return 0.0

    # baseline "Daily Usage Unit" group: Load = Grid = daily usage, Solar = 0
    baseline = n(data.get("chart_daily_usage"))
    if not baseline and chart_opts:
        first = chart_opts[0]
        baseline = n(first.get("grid_units")) + n(first.get("solar_units"))

    groups = [{"label": "DAILY USAGE UNIT", "load": baseline, "grid": baseline, "solar": 0}]
    for o in chart_opts:
        grid = n(o.get("grid_units"))
        solar = n(o.get("solar_units"))
        groups.append({
            "label": (o.get("title") or "Option").upper(),
            "load": grid + solar,
            "grid": grid,
            "solar": solar,
        })

    slide = T.add_slide(prs, page=page, company_name=company_name)
    top = T.add_title(slide, prs, "Daily Power Usage Comparison (Battery + EPC)")
    png = render_usage_chart(groups)
    T.add_image_contain(slide, png, Inches(0.5), top, prs.slide_width - Inches(1.0),
                        prs.slide_height - top - Inches(0.7))


def _placeholder(slide, left, top, w, h):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = T.PANEL
    box.line.color.rgb = T.LINE
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); T._apply_run(r, "No image added", size=13, bold=False, italic=True, color=T.MUTED, font=T.FONT)
